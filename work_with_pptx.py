import os
import pandas as pd

from pptx import Presentation
from pptx.util import Inches

from work_with_exel import get_data_from_sheet



class Present:
    table_top = Inches(0.35)
    table_width = Inches(4.0)
    table_height = Inches(0.8)
    table_lefts = [Inches(2.95), Inches(3) + table_width, Inches(3.05) + 2 * table_width]

    def __init__(self, name_of_file: str, count_slides: int, template_file: str):
        self.name_of_file = name_of_file
        self.count_slides = count_slides
        self.prs = Presentation(template_file)

    def add_slide(self, data: dict):
        images = data["images"]
        tables = data["tables"]
        texts = data["texts"]
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])

        for table in tables:
            num_rows, num_columns = table.shape
            slide.shapes.add_table(num_rows, num_columns, Inches(2), Inches(2), Inches(4), Inches(1.5))

    def save_titul_slide(self, group: str):
        """
        Меняет название группы на титульном слайде шаблона
        """
        shapes = self.prs.slides[0].shapes
        for shape in shapes:
            # Check if the shape is a text box
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if '{group}' in text_frame.text:
                    text_frame.text = text_frame.text.replace("{group}", group)

    def add_table_to_slide(self, data: pd.DataFrame(), slide_index: int):
        """
        Добавляет подготовленную таблицу на слайд шаблонной презы
        :param data: подготволенная таблица из трех столбцов
        :param slide_index: номер слайда, считая от нуля
        """

        shapes = self.prs.slides[slide_index].shapes
        rows, columns = data.shape

        # определяем количество таблиц на листе
        count_of_tables = rows // 16
        if rows - count_of_tables * 16 > 0:
            count_of_tables += 1

        tables = []
        names = data.columns
        for table_i in range(0, count_of_tables):
            tables.append(shapes.add_table(17, columns, self.table_lefts[table_i], self.table_top, self.table_width,
                                           self.table_height).table)
            for i in range(0, len(names)):
                tables[table_i].cell(0, i).text = names[i]
        for row in data.itertuples():
            for index in range(data.columns.size):
                value = str(data.iloc[row.Index, index])
                tables[row.Index // 16].cell(int(row.Index % 16 + 1), int(index)).text = value
        #

    def save(self):
        self.prs.save(self.name_of_file)




def create_new_presentation(name_new_file: str) -> None:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save(f"./pptx files/{name_new_file}.pptx")

